#!/usr/bin/env python
#
# This program is free software; you can redistribute it and/or modify it
# under the terms of the GNU General Public License as published by the
# Free Software Foundation; either version 2 of the License, or (at your
# option) any later version.  See http://www.gnu.org/copyleft/gpl.html for
# the full text of the license.
#
#  https://github.com/python-bugzilla/python-bugzilla/tree/master/examples
#
#  you need to hand install python-bugzilla 1.2.2
#
#
# TODO:
# - coloring when not on track (by expected input spec version 0.3,0.4,... )
# - update to graph to indicate which links are missing
# - make an loop for retrieving all jira items
#   and do an search for where bugzilla id has an value, to limit time.
#
# FIXES:
# - hyperlinks to Jira
# - hyperlinks to bugzilla
# - cross reference to OCF Specs-test Requirements if depends_on field
# - fix for OCF Specs-test Requirements (double links + getting status)
# v5
# - reference to bugs in bugzilla for CTT
# v6
# - added component column
# - added bug state and status for spec CR
# v7
# - added link to bug URL (which might point at word document in Kavi)
#   including check on /latest
# - add check in comments on  [No Test Impact] ==> result in Jira N/A
# v8
# - add check in comments on  [Test Impact] ==> result in Jira should be there
#   (to cancel the [No Test Impact])
# - improved checking on links to kavi (e.g. error only if the url starts with the kavi URL.)
# v9
# - add append to ppt (e.g. new copy to be created)str(args.in_ppt)
# - add sizing of table according slide width
# - date added to generated ppt
# - fix for CTT (rename of component in bugzilla)
# - add option for adding N/A to open source and test requirements individually
#   description on generated info slide
# v10
# - add option to include editorial bugs, italic and includes "(editorial)" in status column
#   also influence the graph (total + not editorial)
# - add option to for adding N/A for CTT individually
# v11 (mark)
# - add option for actions
# - page text in "legend" slide (not perfect though..)
# - added end slide
# v12
# - increased the number of items returned by jira to 1000
#   better solution: make a loop
# v13
# - jira issues in a loop, with filter if bugzilla number is not empty.
# v14 (mark)
# - For JIRA tickets text reads 'Missing' with status of 'unknown' when none are
#   found and OS impact is not tagged.
# v15
# - color coding for priority High and Highest
# - always rename the pptx output file (e.g. add "-out" to the filename)
# v16
# - fix jira: puts in jira tickets in the ppt again
# v17
# - fix jira: multiple bugzilla entries in 1 jira ticket are now conveyed in the ppt.

import time
import os
import random
import sys
import argparse
import traceback
from datetime import datetime
from time import gmtime, strftime
import pprint
import traceback

try:
    import bugzilla
except:
    print("missing bugzilla:")
    print ("Trying to Install required module: bugzilla (should be version 1.2.2 or higher)")
    os.system('python3 -m pip install bugzilla')
import bugzilla


try:
    from jira import JIRA
except:
    print("missing jira  :")
    print ("Trying to Install required module: jira   (should be version 1.2.2 or higher)")
    os.system('python3 -m pip install jira  ')
from jira import JIRA

try:
    import numpy
except:
    print("missing numpy:")
    print ("Trying to Install required module: numpy")
    os.system('python3 -m pip install numpy')
import numpy as np

try:
    import matplotlib.pyplot as plt
except:
    print("missing matplotlib:")
    print ("Trying to Install required module: matplotlib")
    os.system('python3 -m pip install matplotlib')
import matplotlib.pyplot as plt

try:
    from pptx import Presentation
except:
    print("missing ptt:")
    print ("Trying to Install required module: pptx")
    os.system('python3 -m pip install python-pptx')

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE

pp = pprint.PrettyPrinter(indent=4)

if sys.version_info < (3, 5):
    raise Exception("ERROR: Python 3.5 or more is required, you are currently running Python %d.%d!" % (sys.version_info[0], sys.version_info[1]))

TOOL_VERSION = "17"

def filter_bugs(bugs):
    """
    filters out bugs on input criteria
    also filters out bug with severity = "editorial"
    :param bugs: list of bugzilla bugs
    :return: list of bugzilla bug minus the filtered out
    """
    counter = 1
    returned_bugs = []
    for bug in bugs:

        add_bug=True
        if args.open == True:
            if bug.status not in ["IN_PROGRESS", "UNCONFIRMED", "CONFIRMED"]:
                add_bug=False
                if args.verbose:
                    print ("ignoring (not open) :",bug.id, bug.status)
        if args.closed == True:
            if bug.status not in ["RESOLVED", "VERIFIED"]:
                add_bug=False
                if args.verbose:
                    print ("ignoring (not open) :",bug.id, bug.status)
        if args.version is not None:
            if bug.version != args.version:
                add_bug=False
                if args.verbose:
                    print ("ignoring (version) :",bug.id, bug.version)
        if args.state is not None:
            #if  float(bug.cf_uctt_build_version) > float(args.state):
            if bug.cf_uctt_build_version in args.state:
                add_bug=False
                if args.verbose:
                    print ("ignoring (state) :",bug.id, bug.cf_uctt_build_version)
        if args.fromdate is not None:
            size_x = len(args.fromdate)
            creation_time = str(bug.creation_time)[:size_x]
            #print (creation_time, " > ", args.fromdate)
            if creation_time < args.fromdate:
                add_bug=False
                if args.verbose:
                    print ("ignoring (fromdate) :",bug.id, bug.creation_time)
        # remove editoral bugs, if the input arguments is not set
        if args.include_editorial is False:
            if bug.severity in ["editorial"]:
                add_bug=False

        #remove action items, if the input argument is not set
        if args.include_action is False:
            if bug.severity in ["action"]:
                add_bug=False

        if add_bug:
            returned_bugs.append(bug)

    return returned_bugs

def list_all_bugs(bugs):
    """
    list all bugs on the console output
    :param bugs: bugzilla bugs
    """
    counter = 1
    for bug in bugs:
        bug_id = bug.id
        bug_sum = str(bug.summary)
        bug_sev = str(bug.severity)
        bug_stat = str(bug.status)
        bug_comp = str(bug.component)
        try:
            print (counter, bug_id, bug_comp, bug_sum, bug_sev, bug_stat)
        except:
            print (counter, bug_id)
        counter = counter + 1

def ppt_bar_graph(prs, bugs, counter, page_counter, per_page):
    """
    create bar graph
    see: http://python-pptx.readthedocs.io/en/latest/user/charts.html
    :param prs: ppt slide deck
    :param bugs: list of bugs that needs to be added to the
    :param counter: not used
    :param page_counter: not used
    :param per_page: not used
    """
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    total_bugs = len(bugs)
    shapes.title.text = 'Count per component (Total #'+ str(total_bugs)+ ')'

    my_dict = {}
    my_comps = {}
    my_vers = {}
    my_editorial = {}
    editorial_count = 0
    for bug in bugs:
        bug_id = bug.id
        bug_comp = str(bug.component)
        bug_ver = str(bug.version)
        bug_sev = str(bug.severity)
        bug_status = str(bug.status)
        bug_state = str(bug.cf_uctt_build_version)

        val = my_comps.get(bug_comp, 0)
        my_comps[bug_comp] = val + 1

        val = my_editorial.get(bug_comp, 0)
        my_editorial[bug_comp] = val
        if bug_sev in ["editorial"]:
           my_editorial[bug_comp] = val + 1
           editorial_count += 1

        mystring = bug_comp  + " # "+ bug_ver  + " # " +  bug_status
        val = my_vers.get(mystring, 0)
        my_vers[mystring] = val + 1
    print (my_comps)
    #print (my_vers)
    print (my_editorial)
    # define chart data ---------------------
    chart_data = ChartData()
    chart_data.categories = []
    total_string = "total ("+str(total_bugs)+")"
    not_editorial = total_bugs - editorial_count
    not_editorial_string = "not editiorial ("+str(not_editorial)+")"
    series_1 = chart_data.add_series(total_string)
    my_cat = []
    for comp in my_comps:
        my_cat.append(comp)
        value = my_comps[comp]
        series_1.add_data_point(value)


    if (args.include_editorial):
        series_2 = chart_data.add_series(not_editorial_string)
        #series_3 = chart_data.add_series('editorial')
        my_ed_cat = []
        index = 0
        for comp in my_editorial:
            my_ed_cat.append(comp)
            e_value = my_editorial[comp]
            total_value = my_comps[comp]
            #series_3.add_data_point(e_value)
            diff_value = total_value - e_value
            series_2.add_data_point(diff_value)


    chart_data.categories = my_cat
    # add chart to slide --------------------
    width = prs.slide_width

    x = Inches(2)
    y = Inches(2)
    cx = width - (2*y)
    cy = Inches(5)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
    chart = graphic_frame.chart
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels

    if (args.include_editorial):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False


    data_labels.font.size = Pt(13)
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_LABEL_POSITION.INSIDE_END


def ppt_page(prs, bugs, counter, page_counter, per_page, components=None, title=None):
    """
    creates page in slide deck and add table with bug contents on it.
    ppt info at
    # http://python-pptx.readthedocs.io/en/latest/user/quickstart.html
    :param prs: ppt slided
    :param bugs: bugs to be put in the ppt slide (all)
    :param counter: counter where we are
    :param page_counter: counter for the amount pages
    :param per_page: counter for the amount of bugs on a page
    :param components: component name column to be added on the slide
    :param title : additional part of the title of the slide
    """
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    do_comp_column = False

    if title is None:
        shapes.title.text = 'CR overview '+ str(page_counter)
    else:
        shapes.title.text = str(title) + ' CR overview '+ str(page_counter)

    rows =  len(bugs)+1
    cols = 11

    left = Inches(0.0)
    top = Inches(1.5)
    width = prs.slide_width
    height = Inches(0.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(1.0)
    table.columns[3].width = Inches(0.6)
    table.columns[4].width = Inches(1.0)
    table.columns[5].width = Inches(0.6)
    table.columns[6].width = Inches(1.0)
    table.columns[7].width = Inches(0.6)
    table.columns[8].width = Inches(1.0)
    table.columns[9].width = Inches(1.0)
    table.columns[10].width = Inches(0.6)

    total_width = 0
    for column in table.columns:
        total_width += column.width

    if total_width > prs.slide_width:
        print ("setting slide with to:",total_width)
        prs.slide_width = total_width
    else:
        factor = prs.slide_width/total_width
        for column in table.columns:
            new_size = column.width * factor
            column.width = int(new_size)


    # write column headings
    table.cell(0, 0).text = 'CR #'
    table.cell(0, 1).text = 'Description'
    table.cell(0, 2).text = 'Status'
    table.cell(0, 3).text = 'Jira'
    table.cell(0, 4).text = 'Status,\nAssignee'
    table.cell(0, 5).text = 'CTR'
    table.cell(0, 6).text = 'Status'
    table.cell(0, 7).text = 'CTT'
    table.cell(0, 8).text = 'Status'
    table.cell(0, 9).text = 'component'
    table.cell(0,10).text = 'Deps'

    counter = 1
    for bug in bugs:
        #print (dir(bug))
        bug_id = str(bug.id)
        bug_url = str(bug.url)
        bug_sum = str(bug.summary)
        bug_sev = str(bug.severity)
        bug_stat = str(bug.status)
        bug_comp = str(bug.component)
        bug_priority = str(bug.priority)
        bug_state = str(bug.cf_uctt_build_version)
        bug_comments = comments = bug.getcomments()
        try:
            #print (counter, bug_id, bug_comp, bug_sum, bug_sev, bug_stat)
            #
            # first column
            p = table.cell(counter, 0).text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(bug_id)
            r.hyperlink.address = "https://bugzilla.upnp.org/show_bug.cgi?id="+str(bug_id)
            if bug_sev in ["editorial"]:
                #r.font.color.rgb = RGBColor(0x00, 0x7F, 0x50)
                r.font.italic = True
            #
            # 2nd column
            p = table.cell(counter, 1).text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(bug_sum)
            if len(bug_url) > 1:
                r.hyperlink.address = bug_url
                #r.font.color.rgb = RGBColor(0x00, 0x7F, 0x50)
                if bug_sev in ["editorial"]:
                    r.font.italic = True
                if "https://workspace.openconnectivity.org" in bug_url:
                    if "latest" not in bug_url:
                        print ("ERROR: URL in bug_id: ", bug_id, " is not pointing at latest:", bug_url)
            #
            # 3rd column
            p =table.cell(counter, 2).text_frame.paragraphs[0]
            r = p.add_run()
            status_text = bug_stat + " " +str(bug_state)
            if bug_sev in ["editorial"]:
                status_text = status_text + " (editiorial)"
            if bug_sev in ["action"]:
                status_text = status_text + " (action)"
            r.text = status_text
            if bug_sev in ["editorial"]:
                #RGBColor(0xFF, 0x7F, 0x50)
                #r.font.color.rgb = RGBColor(0x00, 0x7F, 0x50)
                r.font.italic = True

            # Open source (Jira)
            # Column 3 & 4
            os_impact_comment_found = check_comments_on_text("[Open Source Impact]",bug_comments)
            no_os_impact_comment_found = check_comments_on_text("[No Open Source Impact]",bug_comments)
            os_impact = True
            #comment_found = check_comments_on_text("JSON format",bug_comments)
            if no_os_impact_comment_found is True:
                os_impact = False
            if os_impact_comment_found is True:
                os_impact = True
            if bug_sev in ["editorial"]:
                os_impact = False
            if os_impact is False:
                table.cell(counter, 3).text = "N/A"
            else:
                jira_number = jira_dict.get(str(bug_id))
                jira_status = "unknown"
                jira_assignee = "unknown"
                print ("bugid ==> jira_number:",str(bug_id), jira_number )
                #if jira_number is not None:
                #    jira_status = get_jira_status(jira_number)
                #    print ("Jira:", jira_number, jira_status)
                p =table.cell(counter, 3).text_frame.paragraphs[0]
                r = p.add_run()
                if jira_number is not None:
                    jira_txt = str(jira_number)
                    # check if it is an csv
                    jira_array = jira_txt.split(",")
                    if len(jira_array) > 1:
                        # do multiple entries.
                        jira_status = ""
                        for jira_entry in jira_array:
                            r.text = str(jira_entry)
                            r.hyperlink.address = "https://jira.iotivity.org/browse/"+jira_entry
                            jira_status += get_jira_status(jira_entry)+ " "
                            r = p.add_run()
                    else:
                        r.text = str(jira_number)
                        r.hyperlink.address = "https://jira.iotivity.org/browse/"+jira_txt
                        jira_status = get_jira_status(jira_number)
                        jira_assignee = get_jira_assignee(jira_number)
                else:
                    r.text = "Missing"
                    #r.text = str(jira_number)
                table.cell(counter, 4).text = str(jira_status+",\n"+jira_assignee)

            # test requirements and CTT
            # column 5 &6
            test_impact_comment_found = check_comments_on_text("[Test Impact]",bug_comments)
            no_test_impact_comment_found = check_comments_on_text("[No Test Impact]",bug_comments)
            test_impact = True
            #comment_found = check_comments_on_text("JSON format",bug_comments)
            if no_test_impact_comment_found is True:
                test_impact = False
            if test_impact_comment_found is True:
                test_impact = True
            if bug_sev in ["editorial"]:
                test_impact = False
            if test_impact is False:
                table.cell(counter, 5).text = "N/A"
            else:
                # test requirement bugs
                test_req_number = test_req_dict.get(str(bug_id))
                test_req_status = get_bugzilla_status(bugs_test_req, test_req_number)
                #table.cell(counter, 5).text = str(test_req_number)
                p = table.cell(counter, 5).text_frame.paragraphs[0]
                if test_req_number is not None:
                    numbers = test_req_number.split(" ")
                    for number in numbers:
                        r = p.add_run()
                        r.text = str(number)
                        r.hyperlink.address = "https://bugzilla.upnp.org/show_bug.cgi?id="+str(number)
                        # add white space
                        r = p.add_run()
                        r.text = " "
                else:
                    r = p.add_run()
                    r.text = str(test_req_number)
                table.cell(counter, 6).text = str(test_req_status)

            # ctt bugs
            # column 7 & 8
            ctt_impact_comment_found = check_comments_on_text("[CTT Impact]",bug_comments)
            no_ctt_impact_comment_found = check_comments_on_text("[No CTT Impact]",bug_comments)
            ctt_number = ctt_dict.get(str(bug_id))
            ctt_status = get_bugzilla_status(bugs_ctt, ctt_number)
            ctt_impact = True
            if no_ctt_impact_comment_found is True:
                ctt_impact = False
            if ctt_impact_comment_found is True:
                ctt_impact = True
            if bug_sev in ["editorial"]:
                ctt_impact = False
            if ctt_impact is False:
                table.cell(counter, 7).text = "N/A"
            else:
                p = table.cell(counter, 7).text_frame.paragraphs[0]
                if ctt_number is not None:
                    numbers = ctt_number.split(" ")
                    for number in numbers:
                        r = p.add_run()
                        r.text = str(number)
                        r.hyperlink.address = "https://bugzilla.upnp.org/show_bug.cgi?id="+str(number)
                        # add white space
                        r = p.add_run()
                        r.text = " "
                else:
                    r = p.add_run()
                    r.text = str(ctt_number)
                table.cell(counter, 8).text = str(ctt_status)

            # bug component it resides in.
            table.cell(counter, 9).text = str(bug_comp)

            # Add dependency information
            if len(bug.depends_on) > 0:
                p = table.cell(counter, 10).text_frame.paragraphs[0]
                for item in bug.depends_on:
                    r = p.add_run()
                    r.text = str(item)
                    r.hyperlink.address = "https://bugzilla.upnp.org/show_bug.cgi?id="+str(item)
                    # add white space
                    r = p.add_run()
                    r.text = " "
            else:
                table.cell(counter, 10).text = "None"

            # make fonts smaller
            # update color for for the priority
            for index in range(cols):
                paragraph = table.cell(counter, index).text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                if bug_priority in ["Highest"]:
                    paragraph.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                if bug_priority in ["High"]:
                    paragraph.font.color.rgb = RGBColor(0x9F, 0x00, 0x00)


        except:
            traceback.print_exc()
        counter = counter + 1

def check_comments_on_text(my_text, comments):
    if comments is not None:
        for comment in comments:
            if my_text in comment['text']:
                #print ("check_comments_on_text: text found!:", my_text, comment['text'])
                return True
    return False

def ppt_info_page(prs):
    """
    add slide deck with "legenda information"

    :param prs: ppt slide deck
    """
    bullet_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'input parameter overview'
    tf = body_shape.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    #p = tf.add_paragraph()
    tf.text = "time generated : " + strftime("%Y-%m-%d %H:%M:%S", gmtime())
    p = tf.add_paragraph()
    p.text = "product        : " + str(args.product)
    p = tf.add_paragraph()
    p.text = "component       : " + str(args.component)
    p = tf.add_paragraph()
    p.text = "open            : " + str(args.open) + "  closed          : " + str(args.closed)
    p = tf.add_paragraph()
    p.text = "version         : " + str(args.version)
    p = tf.add_paragraph()
    p.text = "state           : " + str(args.state)
    p = tf.add_paragraph()
    p.text = "fromdate        : " + str(args.fromdate)
    #p = tf.add_paragraph()
    #p.text = "no_cert        : " + str(args.no_cert)
    p = tf.add_paragraph()
    p.text = "tool version   : " + str(TOOL_VERSION)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "jira status explanation"
    r.hyperlink.address = "https://jira.iotivity.org/plugins/servlet/workflow/thumbnail/getThumbnail?workflowName=Iot+Workflow&stepId=3&width=full&height=full"
    #p.text = "jira status   : " + str("https://jira.iotivity.org/plugins/servlet/workflow/thumbnail/getThumbnail?workflowName=Iot+Workflow&stepId=3&width=full&height=full")
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "bugzilla status explanation"
    r.hyperlink.address = "https://bugzilla.readthedocs.io/en/5.0/using/editing.html#life-cycle-of-a-bug"
    #p.text = "bugzilla bug life cycle info   : " + str("https://bugzilla.readthedocs.io/en/5.0/using/editing.html#life-cycle-of-a-bug")


    #test_impact_comment_found = check_comments_on_text("[Test Impact]",bug_comments)
    #no_test_impact_comment_found = check_comments_on_text("[No Test Impact]",bug_comments)
    p = tf.add_paragraph()
    p.text = 'adding N/A to test requirements: add an comment containing "[No Test Impact]"'
    p = tf.add_paragraph()
    p.text = 'cancelling an N/A to test requirements: add an additional comment containing "[Test Impact]"'
    p.level = 2

    #os_impact_comment_found = check_comments_on_text("[Open Source Impact]",bug_comments)
    #no_os_impact_comment_found = check_comments_on_text("[No Open Source Impact]",bug_comments)
    p = tf.add_paragraph()
    p.text = 'adding N/A to open source: add an comment containing "[No Open Source Impact]"'
    p = tf.add_paragraph()
    p.text = 'cancelling an N/A to open source: add an additional comment containing "[Open Source Impact]"'
    p.level = 2

    #ctt_impact_comment_found = check_comments_on_text("[CTT Impact]",bug_comments)
    #no_ctt_impact_comment_found = check_comments_on_text("[CTT Impact]",bug_comments)
    p = tf.add_paragraph()
    p.text = 'adding N/A to CTT: add an comment containing "[No CTT Impact]"'
    p = tf.add_paragraph()
    p.text = 'cancelling an N/A to CTT: add an additional comment containing "[CTT Impact]"'
    p.level = 2

def ppt_closing_page(prs):
    """
    add closing slide to generated deck

    :param prs: ppt slide deck
    """
    final_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(final_slide_layout)

def list_all_bugs_in_ppt(bugs, filename, components=None, title=None):
    """
    creates the slides
    - bar graph
    - paginated table across slides
    :param bugs: list of bugs to be handled
    :param filename: output filename of the slide deck
    :param components: components to be handled
    :param title : additional part of the title of the slide
    """
    if args.in_ppt is not None:
        ppt_name = args.in_ppt + ".pptx"
        prs = Presentation(pptx=ppt_name)
    else:
        prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]

    counter = 1
    page_counter = 1
    per_page = 8
    per_page_counter = 0
    page_bugs = []
    total_bugs = len(bugs)
    counter_down = len(bugs)

    #ppt_line_graph(prs, bugs, counter, page_counter, per_page)

    ppt_bar_graph(prs, bugs, counter, page_counter, per_page)
    for bug in bugs:
        page_bugs.append(bug)
        #print (counter, len(page_bugs))
        per_page_counter = per_page_counter + 1
        counter_down = counter_down - 1

        print_page = False
        if per_page_counter  == per_page:
            print_page = True;
        if counter_down == 0:
            print_page = True;

        if print_page:
            per_page_counter = 0
            try:
                try:
                    ppt_page(prs, page_bugs, counter, page_counter, per_page, components=components, title=title)
                except:
                    traceback.print_exc()
                page_bugs[:] = []
                page_counter = page_counter +1
            except:
                traceback.print_exc()
        counter = counter + 1
    ppt_info_page(prs)
    ppt_closing_page(prs)
    prs.save(filename +'.pptx')



def bugs_graph(bugs, product=None):
    """
    math plot graph (interactive)
    see:
    http://matplotlib.org/examples/pylab_examples/bar_stacked.html
    :param bugs: bugs to be put in graph
    :param product: product name to be used in the title
    """
    my_dict = {}
    for bug in bugs:
        bug_id = bug.id
        bug_comp = str(bug.component)
        bug_ver = str(bug.version)
        bug_sev = str(bug.severity)
        bug_status = str(bug.status)
        bug_state = str(bug.cf_uctt_build_version)

        mystring = bug_comp + "\n" + bug_ver + "\n" + bug_sev + " " + bug_state +"\n" + bug_status
        val = my_dict.get(mystring, 0)
        my_dict[mystring] = val + 1


    width = 0.35       # the width of the bars: can also be len(x) sequence
    ind = np.arange(len(my_dict))    # the x locations for the groups
    p1 = plt.bar(ind, my_dict.values(), width, color='r')
    plt.ylabel('# bugs')
    if product is None:
        plt.title('bugs per group and category')
    else:
        plt.title('bugs per category for '+ str(product))
    plt.xticks(ind + width/2., my_dict.keys(),fontsize=7)

    now = datetime.strftime(datetime.now(), '%Y-%m-%d %H:%M:%S')

    plt.annotate(now, xy=(.800, .975), xycoords='figure fraction',
                horizontalalignment='left', verticalalignment='top',
                fontsize=10)

    plt.show()


def uprint(*objects, sep=' ', end='\n', file=sys.stdout):
    """
    print that handles utf-8
    :param objects: print arguments
    :param sep: seperator between the commands
    :param end: end of line
    :param file: output destination
    :return:
    """
    enc = file.encoding
    if enc == 'UTF-8':
        print(*objects, sep=sep, end=end, file=file)
    else:
        f = lambda obj: str(obj).encode(enc, errors='backslashreplace').decode(enc)
        print(*map(f, objects), sep=sep, end=end, file=file)


def get_jira_status(issue_number):
    """
    retrieves the status of an jira bug
    :param issue_number: full jira issue number (including prefix)
    :return: status of the bug (open,close, inprogress...)
    """
    my_issue = jira.issue(issue_number)
    #print (my_issue)
    return str(my_issue.fields.status)

def get_jira_assignee(issue_number):
    """
    retrieves the assignee of a jira bug
    :param issue_number: full jira issue number (including prefix)
    :return: assignee of the bug
    """
    my_issue = jira.issue(issue_number)
    #print (my_issue)
    return str(my_issue.fields.assignee)

def get_bugzilla_status(my_bug_list, issue_number):
    """
    retrieves the status of an bugzilla bug
    :param my_bug_list: buglist to search the bug number on
    :param issue_number: bugzilla bugnumber
    :return: status of the bug
    """
    return_string = "unknown."
    if issue_number is None:
        return return_string
    try:
        numbers = issue_number.split(" ")
        return_string = ""
        for number in numbers:
            # bug in accessing directly
            #my_bug = bzapi.getbug(issue_number)
            for my_bug in my_bug_list:
                if str(my_bug.id) == number:
                    return_string = return_string + str(my_bug.status) + " "
    except:
        traceback.print_exc()
    return return_string


def get_bugs_from_bugzilla(bzapi, product, component):
    """
    function to retrieve an bug list per producte/component
    :param bzapi: authenticated bugzilla handle
    :param product: product to retrieve
    :param component: component to retrieve (can be none for all components of the product)
    :return:
    """
    query = bzapi.build_query(
         product=product,
         component=component)
    query["include_fields"] = ["component",
                                "version",
                                "id",
                                "summary",
                                "status",
                                "severity" ,
                                "assigned_to" ,
                                "depends_on" ,
                                "url" ,
                                "comments",
                                "priority",
                                "cf_fix_build",
                                "cf_affected_area",
                                "cf_uctt_build_version",
                                "creation_time",
                                "last_change_time"]
    # query() is what actually performs the query. it's a wrapper around Bug.search
    t1 = time.time()
    my_bugs = bzapi.query(query)
    t2 = time.time()
    print("Found %d bugs with query" % len(my_bugs), product, "  ", component)
    print("Query processing time: %s" % (t2 - t1))
    return my_bugs


def buglist_to_depends_on_dict(my_bug_list, my_dict):
    """
    fills an dict with key-value as depends-on bugnumber-referencednumber
    :param my_bug_list: buglist of product/components
    :param my_dict: dict to build
    """
    for bug in my_bug_list:
        if args.verbose:
            print("Fetched bug #%s:" % bug.id)
            print("  Status    = %s" % bug.status)
            print("  depends_on=", str(bug.depends_on))
        for item in bug.depends_on:
            value = str(bug.id)
            try:
                value = str(bug.id) + " " + str(my_dict[str(item)])
            except KeyError:
                pass
            my_dict[str(item)] = value


def add_key_to_dict(mydict, jira_ticket_id, bugzilla_id):
    key = bugzilla_id
    value = jira_ticket_id
    try:
        # key exist.
        myval = mydict[key]
        new_value = myval+","+value
        mydict[key] = new_value
    except KeyError:
        mydict[key] = value

#
#   main of script
#
print ("*** OCF reports bugzilla/jira (v"+TOOL_VERSION+") ***")

parser = argparse.ArgumentParser()

parser.add_argument( "-ver"     , "--verbose"     , help="Execute in verbose mode", action='store_true')
parser.add_argument( "-uid"     , "--userid"      , help="userid.", nargs='?', const="", required=True)
parser.add_argument( "-pwd"     , "--password"    , help="password.", nargs='?', const="", required=True)

parser.add_argument( "-p"       , "--product"     , default="OCF Specs", help="specify individual products",  nargs='?', const="", required=False)
parser.add_argument( "-c"       , "--component"   , default=None, help='component, 1 omitt for all, multiple: "A" "B" "C"', nargs='*', required=False)
parser.add_argument( "-l"       , "--list"        , help="list products, components that are available", action='store_true')
parser.add_argument( "-open"    , "--open"        , help="cr open e.g. in [IN_PROGRESS, UNCONFIRMED, CONFIRMED]", action='store_true')
parser.add_argument( "-closed"  , "--closed"      , help="cr closed [closed, rejected]", action='store_true')

parser.add_argument( "-include_editorial" , "--include_editorial" , help="include editorial bugs, default editorial bugs are excluded", action='store_true')
parser.add_argument( "-include_action" , "--include_action" , help="include action items, default action items are excluded", action='store_true')
#parser.add_argument( "-no_cert" , "-no_cert"      , help="no certification documents in overview", action='store_true')

parser.add_argument( "-v"       , "--version"     , default=None, help="version [OIC 1.1, OIC 1.0]",  nargs='?', const="", required=False)
parser.add_argument( "-s"       , "--state"       , default=None, help="state [0.3 0.4 0.5 ]",  nargs='?', const="", required=False)
#parser.add_argument( "-impact" , "--impact"      , default=None, help="add impact columns to ppt",  nargs='?', const="", required=False)
#parser.add_argument( "-lp_comp" , "--lp_comp"    , default=None, help="list component in ppt tables",  nargs='?', const="", required=False)

parser.add_argument( "-ppt"     , "--powerpoint"  , default=None, help="ppt filename (without extention)",  nargs='?', const="", required=False)
parser.add_argument( "-in_ppt"  , "--in_ppt"  , default=None, help="ppt filename (in)",  nargs='?', const="", required=False)
parser.add_argument( "-ppt_out_date"  , "--ppt_out_date"  , help="add date to generated ppt output file yyyy-mm-dd-hh-mm-ss", action='store_true')
parser.add_argument( "-f"       , "--fromdate"    , default=None, help="from date  in format 20160805T14:59:13 yyyymmddThh:mm:ss", nargs='?', const="", required=False)

parser.add_argument( "-t"       , "--title"    , default=None, help="additional title on the table pabe", nargs='?', const="", required=False)

args = parser.parse_args()

print("userid            : " + str(args.userid))
print("passwrd           : " + "XXX")

print("product           : " + str(args.product))
print("component         : " + str(args.component))
print("open              : " + str(args.open))
print("closed            : " + str(args.closed))
print("include editorial : " + str(args.include_editorial))
print("include action    : " + str(args.include_action))
print("list              : " + str(args.list))
print("title             : " + str(args.title))
#print("no_cert         : " + str(args.no_cert))

print("powerpoint (in)   : " + str(args.in_ppt))
#print("powerpoint(out)   : " + str(args.powerpoint))
print("powerpoint(date)  : " + str(args.ppt_out_date))
#print("impact          : " + str(args.impact))
#print("lp_comp         : " + str(args.lp_comp))

print("version           : " + str(args.version))
print("state             : " + str(args.state))
print("fromdate          : " + str(args.fromdate))

#print("tool version    : " + TOOL_VERSION)

URL = "bugzilla.openconnectivity.org"
bzapi = bugzilla.Bugzilla(URL)


#
# logging in
#
loggedin=False
if args.userid:
    if args.password:
        #
        # with command line arguments
        #
        bzapi.login(user=args.userid,password=args.password )
        loggedin=True
if not loggedin:
    #
    # log in interactively.
    #
    print ("Bugzilla userid:")
    bzapi.interactive_login()

    print ("loging succesfull")
    print ("")


if args.list:
    # list all components, and then quit
    products = bzapi.getproducts()
    for product in products:
        if "(private)" not in str(product['name']):
            print( "product (-p)    : ", str(product['name']))
            print( "description     : ", str(product['description']))
            comps = bzapi.getcomponents(product['name'])
            if comps is not None:
                #       description     :
                print( "components (-c) :")
            for comp in comps:
                print ( "    ", str(comp))
    exit(0)

#
# JIRA
#
jira = JIRA('https://jira.iotivity.org')
projects = jira.projects()

resp=jira.fields()
fmap = { }
#fr i in resp:
# print ( i.id, id.name)
#pp.pprint (resp)
for i in resp:
     field_name=i[u'name'].encode('ascii','ignore')
     field_id=i[u'id'].encode('ascii','ignore')
     fmap[field_name]=field_id
#if args.verbose:
#    pprint.pprint(fmap)

#print (projects)
#allfields=jira.fields()
#nameMap = {field['name']:field['id'] for field in allfields}

#
# create jira dict
#
jira_dict = {}
request_size=500
got = request_size
total = 0
while got==request_size:
    #issues_in_proj = jira.search_issues('project=IOT and cf[10500] is not EMPTY', startAt = total, maxResults= request_size)
    issues_in_proj_iot = jira.search_issues('project=IOT and cf[10500] is not EMPTY', startAt = total, maxResults= request_size)
    issues_in_proj_lite = jira.search_issues('project=LITE and cf[10500] is not EMPTY', startAt = total, maxResults= request_size)
    issues_in_proj = issues_in_proj_iot + issues_in_proj_lite

    got = len(issues_in_proj)
    total += got

    print ("jira count", len(issues_in_proj))
    for issue_key in issues_in_proj:
        if args.verbose:
            print ( "jira issue number:", issue_key )
        try:
            key = str(issue_key)
            my_issue = jira.issue(key)
            if args.verbose:
                print (my_issue)
                #print ("dir:",dir(my_issue))

            bugzilla_id = my_issue.fields.customfield_10500
            if args.verbose:
                print ( "Bugzilla id", bugzilla_id)
            if bugzilla_id is not None:
                bugzilla_id_str = str(bugzilla_id)
                bugs = bugzilla_id_str.split(",")
                if len(bugs) > 1:
                    for bug in bugs:
                        add_key_to_dict(jira_dict, key, bug)
                else:
                    add_key_to_dict(jira_dict, key, bugzilla_id_str)
        except:
            #traceback.print_exc()
            print ("something went wrong!!")
            traceback.print_exc()

if args.verbose:
#if True:
    print ("jira dictionary linking bugzilla id to jira id:")
    print (jira_dict)
    #exit(0)


#exit(0)
# bugzilla
# do the query for the input args
#
bugs = get_bugs_from_bugzilla(bzapi,args.product, args.component )

#
# do the query in for the test requirements
#
bugs_test_req = get_bugs_from_bugzilla(bzapi,"OCF Specs", "Test requirements" )
test_req_dict = {}
buglist_to_depends_on_dict(bugs_test_req, test_req_dict)
if args.verbose:
    print ("test requirement depends on dictionary:")
    print (test_req_dict)
#
# do the query in for the CTT
#
bugs_ctt = get_bugs_from_bugzilla(bzapi,"OCF (public)", "CTT" )
ctt_dict = {}
buglist_to_depends_on_dict(bugs_ctt, ctt_dict)
if args.verbose:
    print ("ctt depends on dictionary:")
    print (ctt_dict)

print ("filtering bugs")
fbugs = filter_bugs(bugs)

print("Found %d bugs after filter"% len(fbugs))

#if args.powerpoint is not None:
filename = str(args.powerpoint)
if args.in_ppt is not None:
    filename = str(args.in_ppt)

if args.ppt_out_date is True :
    now = time.strftime("-%Y-%m-%d-%H-%M-%S")
    filename = filename + now + ".pptx"
else:
    filename = filename + "-out"

print ("outfilename = ", filename)

#print("powerpoint (in) : " + str(args.in_ppt))
#print("powerpoint(out) : " + str(args.powerpoint))
#print("powerpoint(date) : " + str(args.ppt_out_date))

list_all_bugs_in_ppt(fbugs, filename, components = args.component, title=args.title)


exit(0)
